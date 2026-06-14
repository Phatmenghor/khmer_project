"""
eMenu Cambodia + Tiffany Furniture — Server Setup Guide (PowerPoint)
Clean light theme, easy-to-copy code blocks.
Run: python3 generate-server-setup-pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ORANGE  = RGBColor(0xE8, 0x6B, 0x00)
DARK    = RGBColor(0x1F, 0x2D, 0x3D)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF7, 0xF8, 0xFA)
MGRAY   = RGBColor(0x6B, 0x72, 0x80)
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
BLUE    = RGBColor(0x25, 0x63, 0xEB)
RED     = RGBColor(0xDC, 0x26, 0x26)
YELLOW  = RGBColor(0xF5, 0x9E, 0x0B)
CODE_BG = RGBColor(0xF1, 0xF5, 0xF9)
CODE_FG = RGBColor(0x0F, 0x17, 0x2A)
BORDER  = RGBColor(0xE2, 0xE8, 0xF0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── helpers ───────────────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None, line=None, line_w=0.75):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h, size=11, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def code(slide, lines, l, t, w, h):
    rect(slide, l, t, w, h, fill=CODE_BG, line=BORDER)
    tb = slide.shapes.add_textbox(
        Inches(l+0.12), Inches(t+0.1), Inches(w-0.24), Inches(h-0.2))
    tf = tb.text_frame; tf.word_wrap = False
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = line
        r.font.size = Pt(8.5); r.font.name = "Courier New"
        r.font.color.rgb = CODE_FG

def badge(slide, text, l, t, color=ORANGE):
    rect(slide, l, t, 1.1, 0.26, fill=color)
    txt(slide, text, l+0.05, t+0.03, 1.0, 0.22,
        size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 0.9, fill=WHITE)
    rect(slide, 0, 0, 0.08, 0.9, fill=ORANGE)
    txt(slide, title, 0.25, 0.08, 12.5, 0.5,
        size=22, bold=True, color=DARK)
    if subtitle:
        txt(slide, subtitle, 0.25, 0.58, 12.5, 0.3,
            size=10, color=MGRAY, italic=True)
    rect(slide, 0, 0.9, 13.33, 0.02, fill=BORDER)
    rect(slide, 0, 7.45, 13.33, 0.05, fill=ORANGE)

def row_table(slide, headers, rows, l, t, col_widths):
    h_top = t
    x = l
    for i, (hdr, w) in enumerate(zip(headers, col_widths)):
        rect(slide, x, h_top, w, 0.32, fill=DARK)
        txt(slide, hdr, x+0.08, h_top+0.07, w-0.1, 0.22,
            size=9, bold=True, color=WHITE)
        x += w
    for ri, row in enumerate(rows):
        x = l
        fill = LGRAY if ri % 2 == 0 else WHITE
        for ci, (cell, w) in enumerate(zip(row, col_widths)):
            rect(slide, x, h_top+0.32+ri*0.38, w, 0.38, fill=fill, line=BORDER)
            c = ORANGE if ci == 0 else DARK
            txt(slide, cell, x+0.08, h_top+0.32+ri*0.38+0.08,
                w-0.1, 0.25, size=9, color=c, bold=(ci==0))
            x += w

def step_item(slide, num, title, desc, t):
    rect(slide, 0.3, t, 0.36, 0.36, fill=ORANGE)
    txt(slide, str(num), 0.3, t+0.06, 0.36, 0.26,
        size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, title, 0.8, t+0.02, 12.0, 0.22, size=10, bold=True, color=DARK)
    txt(slide, desc,  0.9, t+0.24, 11.9, 0.22, size=9,  color=MGRAY)
    return t + 0.52


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
rect(s, 0, 0, 13.33, 0.06, fill=ORANGE)
rect(s, 0, 7.44, 13.33, 0.06, fill=ORANGE)
rect(s, 0.4, 1.8, 0.1, 3.5, fill=ORANGE)
txt(s, "Server Setup Guide", 0.7, 1.9, 12, 0.8,
    size=42, bold=True, color=DARK)
txt(s, "eMenu Cambodia  +  Tiffany Furniture KH", 0.7, 2.75, 12, 0.55,
    size=22, color=ORANGE, bold=True)
txt(s, "Complete deployment & configuration reference for both projects",
    0.7, 3.4, 11, 0.4, size=13, color=MGRAY)
txt(s, "Server IP: 167.172.88.194   •   OS: Ubuntu Linux   •   Last Updated: June 2026",
    0.7, 3.9, 11, 0.35, size=11, color=MGRAY)

rect(s, 0.7, 4.6, 3.6, 0.6, fill=LGRAY, line=BORDER)
txt(s, "emenu-cambodia.com", 0.85, 4.7, 3.3, 0.35, size=11, bold=True, color=ORANGE)

rect(s, 4.6, 4.6, 3.8, 0.6, fill=LGRAY, line=BORDER)
txt(s, "tiffanyfurniturekh.com", 4.75, 4.7, 3.5, 0.35, size=11, bold=True, color=BLUE)

rect(s, 8.7, 4.6, 3.9, 0.6, fill=LGRAY, line=BORDER)
txt(s, "167.172.88.194 (VPS)", 8.85, 4.7, 3.6, 0.35, size=11, bold=True, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Architecture Overview
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Architecture Overview", "All services run in Docker on a single Linux VPS — ports bound to 127.0.0.1 (Nginx only)")

row_table(s,
    ["Container", "Image / Port", "Domain", "Project"],
    [
        ["emenu_backend",         "emenu-cambodia/backend:1.0.0  :7070",  "api.emenu-cambodia.com",       "eMenu"],
        ["emenu_owner_frontend",  "owner-frontend  :3000",                "emenu-cambodia.com",            "eMenu"],
        ["emenu_client_frontend", "client-frontend  :3001",               "*.emenu-cambodia.com",          "eMenu"],
        ["tiffany_backend",       "tiffany/backend:1.0.0  :6060",         "api.tiffanyfurniturekh.com",    "Tiffany"],
        ["tiffany_frontend",      "tiffany/frontend:1.0.0  :3002",        "tiffanyfurniturekh.com",        "Tiffany"],
        ["emenu-redis",           "redis:7-alpine  :6379",                "Internal only",                 "Shared"],
        ["postgres17",            "postgres:17  :5432",                   "Internal only",                 "Shared"],
    ],
    0.3, 1.05, [3.0, 3.3, 3.5, 1.9]
)

rect(s, 0.3, 4.5, 12.7, 0.5, fill=LGRAY, line=BORDER)
txt(s, "Security: All app ports bound to 127.0.0.1 (not 0.0.0.0) — only Nginx can reach them from outside.",
    0.5, 4.58, 12.3, 0.35, size=10, color=DARK, bold=True)

txt(s, "SSL Certificates (auto-renew via acme.sh cron daily 11:44)", 0.3, 5.15, 12.7, 0.3,
    size=11, bold=True, color=DARK)
row_table(s,
    ["Domain", "CA", "Cert Path", "Renews"],
    [
        ["emenu-cambodia.com + *.emenu-cambodia.com", "ZeroSSL",       "/etc/ssl/emenu-cambodia/fullchain.pem",    "2026-08-23"],
        ["tiffanyfurniturekh.com + www + api",        "Let's Encrypt", "/etc/ssl/tiffany-furniture/fullchain.pem", "2026-08-14"],
    ],
    0.3, 5.45, [4.2, 2.0, 4.8, 1.7]
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Directory Structure
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Server Directory Structure")

code(s, [
    "/opt/apps/",
    "├── emenu-backend/                   ← eMenu Spring Boot deploy",
    "│   ├── e_menu_cambodia.jar",
    "│   ├── Dockerfile",
    "│   ├── docker-compose.yml",
    "│   ├── logs/",
    "│   └── README.md",
    "├── owner-frontend/                  ← eMenu Owner frontend deploy",
    "│   ├── menu-scanner-frontend-owner/ (source copy)",
    "│   ├── Dockerfile",
    "│   └── docker-compose.yml",
    "├── client-frontend/                 ← eMenu Client frontend deploy",
    "│   ├── menu-scanner-frontend-client/ (source copy)",
    "│   ├── Dockerfile",
    "│   └── docker-compose.yml",
    "├── tiffany-backend/                 ← Tiffany Spring Boot deploy",
    "│   ├── tiffany_backend.jar",
    "│   ├── Dockerfile",
    "│   ├── docker-compose.yml",
    "│   └── README.md",
    "├── tiffany-frontend/                ← Tiffany frontend deploy",
    "│   ├── tiffany-frontend/ (source copy)",
    "│   ├── Dockerfile",
    "│   └── docker-compose.yml",
    "├── tiffany_wong/                    ← Tiffany source code (git repo)",
    "│   ├── tiffany-backend/",
    "│   └── tiffany-frontend/",
    "└── khmer_project/                   ← eMenu source code (git repo)",
], 0.3, 1.0, 6.2, 6.1)

code(s, [
    "/etc/nginx/conf.d/",
    "├── 000-default.conf    ← block direct IP access (return 444)",
    "├── emenu-owner.conf    ← emenu-cambodia.com        → :3000",
    "├── emenu-api.conf      ← api.emenu-cambodia.com    → :7070",
    "├── emenu-client.conf   ← *.emenu-cambodia.com      → :3001",
    "├── tiffany-frontend.conf ← tiffanyfurniturekh.com  → :3002",
    "└── tiffany-api.conf    ← api.tiffanyfurniturekh.com → :6060",
    "",
    "/etc/ssl/",
    "├── emenu-cambodia/",
    "│   ├── fullchain.pem",
    "│   └── key.pem",
    "└── tiffany-furniture/",
    "    ├── fullchain.pem",
    "    └── key.pem",
    "",
    "/opt/tools/",
    "├── redis/docker-compose.yml",
    "└── postgres/docker-compose.yml",
], 6.7, 1.0, 6.3, 6.1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SSL Setup
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "SSL Certificate Setup", "acme.sh + GoDaddy DNS API — wildcard and multi-domain certs")

txt(s, "Get GoDaddy API Key:", 0.3, 1.05, 6.0, 0.3, size=11, bold=True, color=DARK)
txt(s, "https://developer.godaddy.com/keys  →  Create New API Key  →  Environment: Production",
    0.3, 1.35, 6.2, 0.3, size=9, color=MGRAY)

code(s, [
    "# Install acme.sh (one-time)",
    'curl https://get.acme.sh | sh -s email=phatmenghor19@gmail.com',
    "",
    "# ── eMenu (wildcard) ────────────────────────────────────",
    'export GD_Key="YOUR_GODADDY_KEY"',
    'export GD_Secret="YOUR_GODADDY_SECRET"',
    "",
    "~/.acme.sh/acme.sh --issue --dns dns_gd \\",
    "  -d emenu-cambodia.com -d '*.emenu-cambodia.com' \\",
    "  --server letsencrypt",
    "",
    "mkdir -p /etc/ssl/emenu-cambodia",
    "~/.acme.sh/acme.sh --install-cert -d emenu-cambodia.com \\",
    "  --fullchain-file /etc/ssl/emenu-cambodia/fullchain.pem \\",
    "  --key-file       /etc/ssl/emenu-cambodia/key.pem \\",
    '  --reloadcmd      "nginx -s reload"',
], 0.3, 1.7, 6.2, 4.1)

code(s, [
    "# ── Tiffany (multi-domain) ──────────────────────────────",
    'export GD_Key="YOUR_GODADDY_KEY"',
    'export GD_Secret="YOUR_GODADDY_SECRET"',
    "",
    "~/.acme.sh/acme.sh --issue --dns dns_gd \\",
    "  -d tiffanyfurniturekh.com \\",
    "  -d www.tiffanyfurniturekh.com \\",
    "  -d api.tiffanyfurniturekh.com \\",
    "  --server letsencrypt",
    "",
    "mkdir -p /etc/ssl/tiffany-furniture",
    "~/.acme.sh/acme.sh --install-cert -d tiffanyfurniturekh.com \\",
    "  --fullchain-file /etc/ssl/tiffany-furniture/fullchain.pem \\",
    "  --key-file       /etc/ssl/tiffany-furniture/key.pem \\",
    '  --reloadcmd      "nginx -s reload"',
    "",
    "# Check renewal schedule",
    "~/.acme.sh/acme.sh --list",
    "crontab -l | grep acme",
], 6.8, 1.05, 6.2, 5.85)

rect(s, 0.3, 5.9, 6.2, 0.45, fill=RGBColor(0xFF,0xF3,0xCD), line=YELLOW)
txt(s, "⚠  Regenerate GoDaddy API key after each use. Never commit to git. SSL valid 90 days, auto-renews at 60 days.",
    0.45, 5.97, 5.9, 0.33, size=9, color=RGBColor(0x92,0x40,0x00))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Nginx Configuration
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Nginx Configuration", "6 config files — one per service + default block to deny IP access")

configs = [
    ("000-default.conf", "Block direct IP access", [
        "server {",
        "  listen 80 default_server;",
        "  listen 443 ssl default_server;",
        "  server_name _;",
        "  ssl_certificate /etc/ssl/emenu-cambodia/fullchain.pem;",
        "  ssl_certificate_key /etc/ssl/emenu-cambodia/key.pem;",
        "  return 444;",
        "}",
    ]),
    ("emenu-owner.conf", "emenu-cambodia.com → :3000", [
        "server {",
        "  listen 80; listen [::]:80;",
        "  server_name emenu-cambodia.com www.emenu-cambodia.com;",
        "  return 301 https://$host$request_uri;",
        "}",
        "server {",
        "  listen 443 ssl; listen [::]:443 ssl;",
        "  server_name emenu-cambodia.com www.emenu-cambodia.com;",
        "  ssl_certificate /etc/ssl/emenu-cambodia/fullchain.pem;",
        "  ssl_certificate_key /etc/ssl/emenu-cambodia/key.pem;",
        "  ssl_protocols TLSv1.2 TLSv1.3;",
        "  client_max_body_size 50M; server_tokens off;",
        "  location / {",
        "    proxy_pass http://127.0.0.1:3000;",
        "    proxy_http_version 1.1;",
        "    proxy_set_header Host $host;",
        "    proxy_set_header X-Real-IP $remote_addr;",
        "  }",
        "}",
    ]),
    ("emenu-api.conf", "api.emenu-cambodia.com → :7070", [
        "# same structure as owner",
        "# server_name api.emenu-cambodia.com;",
        "# proxy_pass http://127.0.0.1:7070;",
    ]),
]
configs2 = [
    ("emenu-client.conf", "*.emenu-cambodia.com → :3001", [
        "# same structure as owner",
        "# server_name *.emenu-cambodia.com;",
        "# proxy_pass http://127.0.0.1:3001;",
    ]),
    ("tiffany-frontend.conf", "tiffanyfurniturekh.com → :3002", [
        "server {",
        "  listen 80; listen [::]:80;",
        "  server_name tiffanyfurniturekh.com www.tiffanyfurniturekh.com;",
        "  return 301 https://$host$request_uri;",
        "}",
        "server {",
        "  listen 443 ssl; listen [::]:443 ssl;",
        "  server_name tiffanyfurniturekh.com www.tiffanyfurniturekh.com;",
        "  ssl_certificate /etc/ssl/tiffany-furniture/fullchain.pem;",
        "  ssl_certificate_key /etc/ssl/tiffany-furniture/key.pem;",
        "  ssl_protocols TLSv1.2 TLSv1.3;",
        "  client_max_body_size 50M; server_tokens off;",
        "  location / {",
        "    proxy_pass http://127.0.0.1:3002;",
        "    proxy_http_version 1.1;",
        "    proxy_set_header Host $host;",
        "  }",
        "}",
    ]),
    ("tiffany-api.conf", "api.tiffanyfurniturekh.com → :6060", [
        "# same structure as tiffany-frontend",
        "# server_name api.tiffanyfurniturekh.com;",
        "# proxy_pass http://127.0.0.1:6060;",
    ]),
]

for i, (fname, label, lines) in enumerate(configs):
    top = 1.05 + i * 2.05
    rect(s, 0.3, top, 6.2, 0.3, fill=DARK)
    txt(s, fname + "  —  " + label, 0.42, top+0.05, 6.0, 0.22,
        size=9, bold=True, color=WHITE)
    code(s, lines, 0.3, top+0.3, 6.2, 1.68)

for i, (fname, label, lines) in enumerate(configs2):
    top = 1.05 + i * 2.05
    rect(s, 6.85, top, 6.2, 0.3, fill=DARK)
    txt(s, fname + "  —  " + label, 6.97, top+0.05, 6.0, 0.22,
        size=9, bold=True, color=WHITE)
    code(s, lines, 6.85, top+0.3, 6.2, 1.68)

code(s, ["nginx -t && nginx -s reload"], 0.3, 7.12, 12.75, 0.28)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — eMenu Backend Deploy
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "eMenu Backend — Deploy & Update", "Spring Boot JAR → /opt/apps/emenu-backend/ → Docker port 7070")

txt(s, "First Time Deploy", 0.3, 1.05, 6.0, 0.3, size=12, bold=True, color=DARK)
code(s, [
    "# 1. Build JAR",
    "cd /opt/apps/khmer_project/menu-scanner-backend",
    "docker run --rm \\",
    "  -v $(pwd):/app -v maven-cache:/root/.m2 -w /app \\",
    "  maven:3.9-eclipse-temurin-21-alpine \\",
    "  mvn clean package -DskipTests",
    "",
    "# 2. Copy JAR",
    "cp target/e_menu_cambodia.jar /opt/apps/emenu-backend/",
    "",
    "# 3. Build & start",
    "cd /opt/apps/emenu-backend",
    "docker compose up -d --build",
    "",
    "# 4. Verify",
    "docker ps",
    "curl http://localhost:7070/actuator/health",
], 0.3, 1.4, 6.2, 4.5)

txt(s, "Update Deploy", 6.8, 1.05, 6.0, 0.3, size=12, bold=True, color=DARK)
code(s, [
    "# 1. Pull latest code",
    "cd /opt/apps/khmer_project",
    "git pull origin main",
    "",
    "# 2. Rebuild JAR",
    "cd menu-scanner-backend",
    "docker run --rm \\",
    "  -v $(pwd):/app -v maven-cache:/root/.m2 -w /app \\",
    "  maven:3.9-eclipse-temurin-21-alpine \\",
    "  mvn clean package -DskipTests",
    "",
    "# 3. Copy new JAR",
    "cp target/e_menu_cambodia.jar /opt/apps/emenu-backend/",
    "",
    "# 4. Rebuild & restart",
    "cd /opt/apps/emenu-backend",
    "docker compose up -d --build",
], 6.8, 1.4, 6.2, 4.5)

code(s, [
    "docker logs -f emenu_backend    # view logs",
    "docker restart emenu_backend    # quick restart (no rebuild)",
], 0.3, 6.05, 12.75, 0.65)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — eMenu Frontends Deploy
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "eMenu Frontends — Deploy & Update", "Owner :3000 (emenu-cambodia.com)  •  Client :3001 (*.emenu-cambodia.com)")

txt(s, "Owner Frontend — /opt/apps/owner-frontend/", 0.3, 1.05, 6.2, 0.3,
    size=11, bold=True, color=DARK)
code(s, [
    "# First deploy",
    "cp -r /opt/apps/khmer_project/menu-scanner-frontend-owner \\",
    "       /opt/apps/owner-frontend/",
    "cd /opt/apps/owner-frontend",
    "docker compose up -d --build",
    "",
    "# Update",
    "cd /opt/apps/khmer_project && git pull origin main",
    "rm -rf /opt/apps/owner-frontend/menu-scanner-frontend-owner",
    "cp -r menu-scanner-frontend-owner /opt/apps/owner-frontend/",
    "cd /opt/apps/owner-frontend",
    "docker compose up -d --build",
], 0.3, 1.4, 6.2, 3.5)

txt(s, "Client Frontend — /opt/apps/client-frontend/", 6.8, 1.05, 6.2, 0.3,
    size=11, bold=True, color=DARK)
code(s, [
    "# First deploy",
    "cp -r /opt/apps/khmer_project/menu-scanner-frontend-client \\",
    "       /opt/apps/client-frontend/",
    "cd /opt/apps/client-frontend",
    "docker compose up -d --build",
    "",
    "# Update",
    "cd /opt/apps/khmer_project && git pull origin main",
    "rm -rf /opt/apps/client-frontend/menu-scanner-frontend-client",
    "cp -r menu-scanner-frontend-client /opt/apps/client-frontend/",
    "cd /opt/apps/client-frontend",
    "docker compose up -d --build",
], 6.8, 1.4, 6.2, 3.5)

txt(s, "Dockerfile pattern (same for both, change port 3000/3001):", 0.3, 5.05, 12.7, 0.3,
    size=11, bold=True, color=DARK)
code(s, [
    "FROM node:20-alpine AS deps",
    "WORKDIR /app",
    "COPY menu-scanner-frontend-owner/package.json menu-scanner-frontend-owner/package-lock.json ./",
    "RUN npm ci",
    "FROM node:20-alpine AS builder",
    "WORKDIR /app",
    "COPY --from=deps /app/node_modules ./node_modules",
    "COPY menu-scanner-frontend-owner/ .",
    "RUN npm run build",
    "FROM node:20-alpine AS runner",
    "WORKDIR /app",
    "ENV NODE_ENV=production PORT=3000 HOSTNAME=0.0.0.0",
    "COPY --from=builder /app/public ./public",
    "COPY --from=builder --chown=nextjs:nodejs /app/.next/standalone ./",
    "COPY --from=builder --chown=nextjs:nodejs /app/.next/static ./.next/static",
    'CMD ["node", "server.js"]',
], 0.3, 5.38, 12.75, 1.97)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Tiffany Backend Deploy
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Tiffany Backend — Deploy & Update", "Spring Boot JAR → /opt/apps/tiffany-backend/ → Docker port 6060")

txt(s, "First Time Deploy", 0.3, 1.05, 6.0, 0.3, size=12, bold=True, color=DARK)
code(s, [
    "# 1. Build JAR",
    "cd /opt/apps/tiffany_wong/tiffany-backend",
    "docker run --rm \\",
    "  -v $(pwd):/app -v maven-cache:/root/.m2 -w /app \\",
    "  maven:3.9-eclipse-temurin-21-alpine \\",
    "  mvn clean package -DskipTests",
    "",
    "# 2. Copy JAR",
    "cp target/tiffany_backend.jar /opt/apps/tiffany-backend/",
    "",
    "# 3. Build & start",
    "cd /opt/apps/tiffany-backend",
    "docker compose up -d --build",
    "",
    "# 4. Verify",
    "docker ps",
    "curl http://localhost:6060/actuator/health",
], 0.3, 1.4, 6.2, 4.5)

txt(s, "Update Deploy", 6.8, 1.05, 6.0, 0.3, size=12, bold=True, color=DARK)
code(s, [
    "# 1. Pull latest code",
    "cd /opt/apps/tiffany_wong",
    "git pull origin main",
    "",
    "# 2. Rebuild JAR",
    "cd tiffany-backend",
    "docker run --rm \\",
    "  -v $(pwd):/app -v maven-cache:/root/.m2 -w /app \\",
    "  maven:3.9-eclipse-temurin-21-alpine \\",
    "  mvn clean package -DskipTests",
    "",
    "# 3. Copy new JAR",
    "cp target/tiffany_backend.jar /opt/apps/tiffany-backend/",
    "",
    "# 4. Rebuild & restart",
    "cd /opt/apps/tiffany-backend",
    "docker compose up -d --build",
], 6.8, 1.4, 6.2, 4.5)

code(s, [
    "docker logs -f tiffany_backend    # view logs",
    "curl https://api.tiffanyfurniturekh.com/actuator/health",
], 0.3, 6.05, 12.75, 0.65)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Tiffany Frontend Deploy
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Tiffany Frontend — Deploy & Update", "Next.js → /opt/apps/tiffany-frontend/ → Docker port 3002")

txt(s, "First Time Deploy", 0.3, 1.05, 6.0, 0.3, size=12, bold=True, color=DARK)
code(s, [
    "# 1. Copy source",
    "cp -r /opt/apps/tiffany_wong/tiffany-frontend \\",
    "       /opt/apps/tiffany-frontend/tiffany-frontend",
    "",
    "# 2. Build & start",
    "cd /opt/apps/tiffany-frontend",
    "docker compose up -d --build",
    "",
    "# 3. Verify",
    "docker ps",
    "curl https://tiffanyfurniturekh.com",
], 0.3, 1.4, 6.2, 3.2)

txt(s, "Update Deploy", 6.8, 1.05, 6.0, 0.3, size=12, bold=True, color=DARK)
code(s, [
    "# 1. Pull latest code",
    "cd /opt/apps/tiffany_wong",
    "git pull origin main",
    "",
    "# 2. Copy updated source",
    "rm -rf /opt/apps/tiffany-frontend/tiffany-frontend",
    "cp -r tiffany-frontend /opt/apps/tiffany-frontend/",
    "",
    "# 3. Rebuild & restart",
    "cd /opt/apps/tiffany-frontend",
    "docker compose up -d --build",
], 6.8, 1.4, 6.2, 3.2)

txt(s, "Dockerfile:", 0.3, 4.75, 12.7, 0.3, size=11, bold=True, color=DARK)
code(s, [
    "FROM node:20-alpine AS deps",
    "WORKDIR /app",
    "COPY tiffany-frontend/package.json tiffany-frontend/package-lock.json ./",
    "RUN npm ci",
    "FROM node:20-alpine AS builder",
    "WORKDIR /app",
    "COPY --from=deps /app/node_modules ./node_modules",
    "COPY tiffany-frontend/ .",
    "RUN npm run build",
    "FROM node:20-alpine AS runner",
    "WORKDIR /app",
    "ENV NODE_ENV=production PORT=3002 HOSTNAME=0.0.0.0",
    "RUN addgroup --system --gid 1001 nodejs && adduser --system --uid 1001 nextjs",
    "COPY --from=builder /app/public ./public",
    "COPY --from=builder --chown=nextjs:nodejs /app/.next/standalone ./",
    "COPY --from=builder --chown=nextjs:nodejs /app/.next/static ./.next/static",
    "USER nextjs",
    'CMD ["node", "server.js"]',
], 0.3, 5.08, 12.75, 2.27)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Database & Init Data
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Database Setup & Production Init Data", "PostgreSQL 17 — Hibernate auto-creates tables on first start")

row_table(s,
    ["Setting", "Value"],
    [
        ["Host",     "167.172.88.194  (or DB_HOST env var)"],
        ["Port",     "5432"],
        ["Database", "e_menu_platform"],
        ["Username", "postgres"],
        ["DDL auto", "update  (Hibernate creates/alters tables on startup)"],
    ],
    0.3, 1.05, [3.0, 9.7]
)

txt(s, "Run scripts/init-production-data.sql on e_menu_platform after first start:", 0.3, 3.0, 12.7, 0.3,
    size=11, bold=True, color=DARK)
code(s, [
    "-- Creates: PLATFORM_OWNER role, CUSTOMER role, platform owner user",
    "--          user profile, role assignment, 3 subscription plans",
    "-- Safe to re-run — all inserts use WHERE NOT EXISTS",
    "",
    "-- Verify:",
    "SELECT name, price, status, duration_type FROM subscription_plans WHERE is_deleted = false;",
    "-- Expected: 3 rows  →  1 Week (WEEKLY) / 1 Month (MONTHLY) / 1 Year (YEARLY)",
], 0.3, 3.35, 12.75, 2.0)

txt(s, "Redis Cache:", 0.3, 5.5, 12.7, 0.3, size=11, bold=True, color=DARK)
code(s, [
    "# Subscription plans cache TTL = 6 hours",
    "# If pricing section shows empty after inserting data via SQL → flush Redis:",
    "docker exec -it <redis-container-id> redis-cli FLUSHDB",
    "",
    "# Cache clears AUTOMATICALLY when plans are updated via the admin API (CacheEvict annotation)",
    "# Manual FLUSHDB only needed when inserting data directly via SQL",
], 0.3, 5.82, 12.75, 1.53)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Location Data Migration
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Location Data Migration", "Province / District / Commune / Village — source DB (bigint PKs) → production (UUID PKs)")

y = step_item(s, 1, "Open pgAdmin → source database → Query Tool", "Run the source database (account_online or similar)", 1.05)
y = step_item(s, 2, "Run Part A queries from  scripts/migrate-location-data.sql", "4 queries: province, district, commune, village — each generates INSERT statements", y)
y = step_item(s, 3, "Download each result as CSV", "Data Output panel → ↓ Download icon → Download as CSV  (province.csv, district.csv, commune.csv, village.csv)", y)
y = step_item(s, 4, "Convert CSV → SQL  (Notepad / text editor)", 'Ctrl+H → find: \"  replace: (empty) → delete header line → Save As .sql (UTF-8 encoding)', y)
y = step_item(s, 5, "Open pgAdmin → e_menu_platform (production) → Query Tool", "Ctrl+O to open file", y)
y = step_item(s, 6, "Run in ORDER — FK dependency chain", "1. province.sql   2. district.sql   3. commune.sql   4. village.sql", y)

code(s, [
    "-- Part C: Verify counts match between source and production",
    "SELECT 'province', COUNT(*) FROM location_province_cbc WHERE is_deleted=false",
    "UNION ALL SELECT 'district', COUNT(*) FROM location_district_cbc WHERE is_deleted=false",
    "UNION ALL SELECT 'commune',  COUNT(*) FROM location_commune_cbc  WHERE is_deleted=false",
    "UNION ALL SELECT 'village',  COUNT(*) FROM location_village_cbc  WHERE is_deleted=false;",
], 0.3, y+0.1, 12.75, 1.1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Telegram Setup
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Telegram Setup", "Webhook + Login Widget (eMenu)  •  Mini App (Tiffany)")

txt(s, "eMenu — Bot Webhook & Login", 0.3, 1.05, 6.2, 0.3, size=12, bold=True, color=DARK)
code(s, [
    "# Set webhook (run once after backend deploy)",
    "curl -s 'https://api.telegram.org/bot<TOKEN>/setWebhook' \\",
    "  --data-urlencode \\",
    "  'url=https://api.emenu-cambodia.com/api/v1/telegram/webhook'",
    "",
    "# Verify",
    "curl 'https://api.telegram.org/bot<TOKEN>/getWebhookInfo'",
    "",
    "# BotFather — enable login widget",
    "# /setdomain → select bot → emenu-cambodia.com",
    "# (covers all subdomains automatically)",
    "",
    "# application-prod.yaml",
    "telegram:",
    "  bot:",
    "    webhook-url: '${TELEGRAM_WEBHOOK_URL:https://api.emenu-cambodia.com}'",
    "    auto-setup: true",
], 0.3, 1.4, 6.2, 4.5)

txt(s, "Tiffany — Mini App Setup", 6.8, 1.05, 6.2, 0.3, size=12, bold=True, color=DARK)
code(s, [
    "# BotFather commands:",
    "# /newapp → select bot",
    "#   Title: Tiffany Furniture KH",
    "#   Web App URL: https://tiffanyfurniturekh.com/",
    "#   Short name: tiffany",
    "",
    "# /setmenubutton → select bot",
    "#   Button text: Open Shop",
    "#   Web App URL: https://tiffanyfurniturekh.com/",
    "",
    "# /setdomain → select bot → tiffanyfurniturekh.com",
    "",
    "# Mini App URL for users:",
    "# t.me/YOUR_BOT_USERNAME/tiffany",
    "",
    "# Add to Next.js layout.tsx:",
    "import Script from 'next/script'",
    "<Script",
    '  src="https://telegram.org/js/telegram-web-app.js"',
    '  strategy="beforeInteractive"',
    "/>",
], 6.8, 1.4, 6.2, 4.5)

code(s, ["curl https://api.emenu-cambodia.com/actuator/health    # verify backend webhook is live"], 0.3, 6.05, 12.75, 0.45)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — DNS Records
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "DNS Records (GoDaddy)", "All A records point to server IP: 167.172.88.194  •  TTL: 600")

txt(s, "emenu-cambodia.com", 0.3, 1.05, 6.2, 0.3, size=12, bold=True, color=ORANGE)
row_table(s,
    ["Type", "Name", "Value", "Purpose"],
    [
        ["A", "@",   "167.172.88.194", "Owner frontend"],
        ["A", "www", "167.172.88.194", "Owner frontend (www)"],
        ["A", "api", "167.172.88.194", "Backend API"],
        ["A", "*",   "167.172.88.194", "Wildcard → client frontend"],
    ],
    0.3, 1.4, [1.0, 1.5, 3.5, 3.95]
)

txt(s, "tiffanyfurniturekh.com", 0.3, 3.8, 6.2, 0.3, size=12, bold=True, color=BLUE)
row_table(s,
    ["Type", "Name", "Value", "Purpose"],
    [
        ["A", "@",   "167.172.88.194", "Frontend"],
        ["A", "www", "167.172.88.194", "Frontend (www)"],
        ["A", "api", "167.172.88.194", "Backend API"],
    ],
    0.3, 4.15, [1.0, 1.5, 3.5, 3.95]
)

txt(s, "Nameservers (GoDaddy):", 6.8, 1.05, 6.0, 0.3, size=12, bold=True, color=DARK)
code(s, [
    "ns47.domaincontrol.com",
    "ns48.domaincontrol.com",
], 6.8, 1.4, 6.2, 0.65)

txt(s, "Verify DNS propagation:", 6.8, 2.2, 6.0, 0.3, size=12, bold=True, color=DARK)
code(s, [
    "# Check A record",
    "nslookup emenu-cambodia.com",
    "nslookup api.tiffanyfurniturekh.com",
    "",
    "# Or use online tool:",
    "# https://dnschecker.org",
], 6.8, 2.55, 6.2, 1.8)

txt(s, "Wait 5–30 min for propagation after adding/changing DNS records.", 6.8, 4.5, 6.2, 0.3,
    size=10, color=MGRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — New Server Checklist
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Moving to a New Server — Complete Checklist")

checklist = [
    ("1. Provision Server",   "Get Linux VPS  •  Install Docker (curl -fsSL https://get.docker.com | sh)  •  Install Nginx + acme.sh"),
    ("2. Update DNS",         "GoDaddy → update all A records to new server IP  •  Wait 5–30 min for propagation"),
    ("3. Issue SSL Certs",    "Run acme.sh for emenu-cambodia.com (wildcard) and tiffanyfurniturekh.com  •  Install to /etc/ssl/"),
    ("4. Nginx",              "Create all 6 conf files in /etc/nginx/conf.d/  •  nginx -t && nginx -s reload"),
    ("5. Start DB & Redis",   "cd /opt/tools/postgres && docker compose up -d  •  cd /opt/tools/redis && docker compose up -d"),
    ("6. Deploy eMenu Backend",     "Build JAR → copy to /opt/apps/emenu-backend/ → docker compose up -d --build"),
    ("7. Deploy eMenu Frontends",   "Copy source → /opt/apps/owner-frontend/ and /opt/apps/client-frontend/ → docker compose up -d --build"),
    ("8. Deploy Tiffany Backend",   "Build JAR → copy to /opt/apps/tiffany-backend/ → docker compose up -d --build"),
    ("9. Deploy Tiffany Frontend",  "Copy source → /opt/apps/tiffany-frontend/tiffany-frontend/ → docker compose up -d --build"),
    ("10. Init Database",     "Run scripts/init-production-data.sql  •  Run scripts/migrate-location-data.sql (Parts A+B)"),
    ("11. Telegram",          "Set webhook URL  •  BotFather /setdomain for both bots  •  /newapp for Tiffany Mini App"),
    ("12. Verify All",        "docker ps  •  curl all health endpoints  •  open all URLs in browser"),
]

for i, (title, desc) in enumerate(checklist):
    col = 0 if i < 6 else 1
    row = i if i < 6 else i - 6
    l = 0.3 + col * 6.55
    t = 1.05 + row * 0.98
    rect(s, l, t, 0.35, 0.35, fill=ORANGE)
    txt(s, title.split(".")[0], l, t+0.06, 0.35, 0.25,
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title.split(". ", 1)[1], l+0.45, t+0.03, 6.0, 0.28,
        size=9, bold=True, color=DARK)
    txt(s, desc, l+0.45, t+0.32, 5.95, 0.6, size=8, color=MGRAY, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Quick Reference
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Quick Reference Commands")

code(s, [
    "# ── View all containers ──────────────────────────────────────",
    "docker ps",
    "",
    "# ── Restart a service ────────────────────────────────────────",
    "docker restart emenu_backend",
    "docker restart tiffany_backend",
    "",
    "# ── View live logs ───────────────────────────────────────────",
    "docker logs -f emenu_backend",
    "docker logs -f tiffany_backend",
    "docker logs -f tiffany_frontend",
    "",
    "# ── Flush Redis cache ────────────────────────────────────────",
    "docker exec -it <redis-container-id> redis-cli FLUSHDB",
    "",
    "# ── Rebuild after code change (eMenu) ────────────────────────",
    "cd /opt/apps/emenu-backend && docker compose up -d --build",
    "cd /opt/apps/owner-frontend && docker compose up -d --build",
    "cd /opt/apps/client-frontend && docker compose up -d --build",
], 0.3, 1.05, 6.2, 6.1)

code(s, [
    "# ── Rebuild after code change (Tiffany) ─────────────────────",
    "cd /opt/apps/tiffany-backend && docker compose up -d --build",
    "cd /opt/apps/tiffany-frontend && docker compose up -d --build",
    "",
    "# ── Check SSL certs & renewal schedule ───────────────────────",
    "~/.acme.sh/acme.sh --list",
    "crontab -l | grep acme",
    "",
    "# ── Force SSL renewal ────────────────────────────────────────",
    "~/.acme.sh/acme.sh --renew -d emenu-cambodia.com --force",
    "~/.acme.sh/acme.sh --renew -d tiffanyfurniturekh.com --force",
    "",
    "# ── Nginx test & reload ──────────────────────────────────────",
    "nginx -t && nginx -s reload",
    "",
    "# ── Health checks ────────────────────────────────────────────",
    "curl https://api.emenu-cambodia.com/actuator/health",
    "curl https://api.tiffanyfurniturekh.com/actuator/health",
    "curl https://emenu-cambodia.com",
    "curl https://tiffanyfurniturekh.com",
], 6.8, 1.05, 6.2, 6.1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Status Checklist
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Setup Complete — All Services Status")

done = [
    (GREEN, "PostgreSQL 17",            "Running on 127.0.0.1:5432  •  Database: e_menu_platform"),
    (GREEN, "Redis 7",                  "Running on 127.0.0.1:6379  •  Cache TTL: 6 hours"),
    (GREEN, "eMenu Backend",            "Running on 127.0.0.1:7070  •  Healthy  •  api.emenu-cambodia.com"),
    (GREEN, "eMenu Owner Frontend",     "Running on 127.0.0.1:3000  •  emenu-cambodia.com"),
    (GREEN, "eMenu Client Frontend",    "Running on 127.0.0.1:3001  •  *.emenu-cambodia.com"),
    (GREEN, "Tiffany Backend",          "Running on 127.0.0.1:6060  •  Healthy  •  api.tiffanyfurniturekh.com"),
    (GREEN, "Tiffany Frontend",         "Running on 127.0.0.1:3002  •  tiffanyfurniturekh.com"),
    (GREEN, "SSL — eMenu",              "*.emenu-cambodia.com  •  ZeroSSL  •  Auto-renew 2026-08-23"),
    (GREEN, "SSL — Tiffany",            "tiffanyfurniturekh.com  •  Let's Encrypt  •  Auto-renew 2026-08-14"),
    (GREEN, "Nginx",                    "6 conf files  •  IP direct access blocked  •  All domains routed"),
    (GREEN, "Telegram — eMenu",         "Webhook active  •  Login domain: emenu-cambodia.com"),
    (GREEN, "Telegram — Tiffany",       "Mini App registered  •  tiffanyfurniturekh.com"),
    (GREEN, "Init Data",                "Platform owner + 3 subscription plans  •  Location data migrated"),
]

for i, (color, service, status) in enumerate(done):
    col = 0 if i < 7 else 1
    row = i if i < 7 else i - 7
    l = 0.3 + col * 6.55
    t = 1.05 + row * 0.84
    rect(s, l, t, 0.36, 0.36, fill=color)
    txt(s, "✓", l, t+0.06, 0.36, 0.26,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, service, l+0.46, t+0.03, 2.8, 0.25, size=10, bold=True, color=DARK)
    txt(s, status,  l+0.46, t+0.28, 5.9, 0.48, size=8, color=MGRAY)


OUTPUT = "/home/user/khmer_project/eMenu_Cambodia_Server_Setup_Guide.pptx"
prs.save(OUTPUT)
print(f"Saved → {OUTPUT}")
